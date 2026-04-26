"""Investor seed deck — 10 slides, 16:9, image-rich."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ASSETS = Path(__file__).parent / "deck_assets"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x29, 0xB6, 0xF6)
PINK = RGBColor(0xFF, 0x3D, 0x7F)
GOLD = RGBColor(0xF7, 0xB7, 0x31)
DARK = RGBColor(0x0E, 0x0E, 0x12)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x5F, 0x6D)


def style_run(run, *, size=18, color=NAVY, bold=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font


def add_text(slide, x, y, w, h, text, *, size=18, color=NAVY, bold=False, align=None, anchor=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    style_run(r, size=size, color=color, bold=bold)
    return tf


def add_filled(slide, x, y, w, h, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_pill(slide, x, y, w, h, text, fill, fg, *, size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.5
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    style_run(r, size=size, color=fg, bold=True)


def add_footer(slide, num, total):
    add_filled(slide, 0, 7.0, 13.333, 0.5, NAVY)
    add_text(slide, 0.4, 7.06, 6, 0.4, "clipsmith · seed round · confidential",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 12.4, 7.06, 0.6, 0.4, f"{num} / {total}",
             size=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def title_block(slide, eyebrow, title):
    # Coloured bar
    add_filled(slide, 0.5, 0.5, 0.15, 0.55, PINK)
    add_text(slide, 0.8, 0.45, 8, 0.4, eyebrow, size=12, color=PINK, bold=True)
    add_text(slide, 0.8, 0.75, 12, 0.7, title, size=30, color=NAVY, bold=True)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 10

    # ============ SLIDE 1 — COVER ============
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(ASSETS / "cover_bg.png"), 0, 0, width=Inches(13.333), height=Inches(7.5))

    # Title block
    add_text(s, 0.8, 2.5, 9, 1.2, "clipsmith", size=88, color=WHITE, bold=True)
    add_filled(s, 0.85, 3.7, 0.6, 0.08, PINK)
    add_text(s, 0.8, 3.85, 11, 0.7, "Where creators build careers, not just followings.",
             size=24, color=WHITE)
    add_text(s, 0.8, 4.55, 11, 0.5, "AI-native short-form video, creator-first economics, professional editing in your pocket.",
             size=14, color=RGBColor(0xCC, 0xD8, 0xE6))

    # Tag pills
    for i, (label, fill) in enumerate([("$3.5M seed", PINK), ("18-month runway", ACCENT), ("iOS launch Q3 2026", GOLD)]):
        add_pill(s, 0.8 + i * 2.0, 5.4, 1.85, 0.4, label, fill, WHITE, size=11)

    # Footer-ish
    add_text(s, 0.8, 6.6, 12, 0.4, "Confidential investor materials · prepared April 2026",
             size=10, color=RGBColor(0x99, 0xA8, 0xC0))

    # ============ SLIDE 2 — PROBLEM ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "01 · The problem", "Short-form video became a billboard for ads, not a career for creators.")

    # Three pain pillars
    pain = [
        ("Opaque payouts",
         "TikTok's Creator Fund pays $0.02–$0.04 per 1,000 views. Most creators have no idea what they will be paid until the cheque clears."),
        ("Toy-grade editing",
         "In-app editors stop at trim + filters. Anyone serious exports to CapCut, Premiere, or Final Cut. The market splits the value."),
        ("No path to revenue",
         "Sub-100k creators can't tip, can't sub, can't sell. Brand-deal matching is a Discord DM, not a marketplace."),
    ]
    for i, (heading, body) in enumerate(pain):
        x = 0.8 + i * 4.15
        add_filled(s, x, 2.0, 3.95, 4.6, WHITE)
        add_filled(s, x, 2.0, 3.95, 0.18, PINK)
        add_text(s, x + 0.3, 2.45, 3.4, 0.6, heading, size=20, color=NAVY, bold=True)
        add_text(s, x + 0.3, 3.2, 3.4, 3.2, body, size=14, color=GRAY)
        add_text(s, x + 0.3, 5.7, 3.4, 0.4, f"#{i+1}", size=72, color=RGBColor(0xE0, 0xE5, 0xEE), bold=True)

    add_footer(s, 2, total)

    # ============ SLIDE 3 — SOLUTION ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "02 · Solution", "An AI-native, creator-first video platform — built end-to-end, not bolted on.")

    # Three columns of solution
    sol = [
        ("Pro editor in 90s",
         "Multi-track timeline, AI captions, scene detection, voice de-noise, BG removal, chroma-key — all on-device or server-side, no export.",
         ACCENT),
        ("Transparent 70/30",
         "Tipping (Stripe), subscriptions (Apple IAP), brand-deals, premium content. Creators see exactly which signal earned them what.",
         PINK),
        ("Discovery you can read",
         "Algorithm exposes the four feature scores per video — interest, virality, freshness, community — so creators can optimise without guessing.",
         GOLD),
    ]
    for i, (heading, body, c) in enumerate(sol):
        x = 0.8 + i * 4.15
        add_filled(s, x, 2.0, 3.95, 4.6, WHITE)
        # icon dot
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(2.4), Inches(0.6), Inches(0.6))
        circle.fill.solid(); circle.fill.fore_color.rgb = c
        circle.line.fill.background()
        add_text(s, x + 1.05, 2.45, 2.6, 0.6, heading, size=18, color=NAVY, bold=True)
        add_text(s, x + 0.3, 3.4, 3.4, 3.0, body, size=13, color=GRAY)

    add_pill(s, 0.8, 6.6, 11.7, 0.4,
             "30k LOC backend · 14k LOC frontend · 90 SQL tables · 498 unit + 10 e2e tests · 5 architectural contracts · all CI green",
             NAVY, WHITE, size=10)
    add_footer(s, 3, total)

    # ============ SLIDE 4 — WHY NOW ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "03 · Why now", "The market just hit the inflection. Bottoms-up, every signal points the same way.")

    s.shapes.add_picture(str(ASSETS / "trend_market.png"), Inches(0.6), Inches(1.8), width=Inches(7.5))
    s.shapes.add_picture(str(ASSETS / "creator_growth.png"), Inches(0.6), Inches(4.6), width=Inches(7.5))

    # Right column callouts
    callouts = [
        ("$540B", "short-form ad spend by 2030", PINK),
        ("545M", "creators monetising by 2030", ACCENT),
        ("17%", "TikTok-share at risk to creator-revolt", GOLD),
        ("9-10mo", "to ship full-stack vs. legacy 24mo", NAVY),
    ]
    for i, (n, sub, c) in enumerate(callouts):
        y = 1.95 + i * 1.25
        add_filled(s, 8.5, y, 4.4, 1.05, WHITE)
        add_filled(s, 8.5, y, 0.12, 1.05, c)
        add_text(s, 8.75, y + 0.05, 4, 0.6, n, size=32, color=c, bold=True)
        add_text(s, 8.75, y + 0.55, 4, 0.5, sub, size=12, color=GRAY)

    add_footer(s, 4, total)

    # ============ SLIDE 5 — PRODUCT ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "04 · Product", "Three surfaces, one codebase — feed, edit, monetise.")

    # Three phone screenshots
    s.shapes.add_picture(str(ASSETS / "phone_feed.png"), Inches(1.0), Inches(1.8), height=Inches(5.0))
    s.shapes.add_picture(str(ASSETS / "phone_editor.png"), Inches(5.2), Inches(1.8), height=Inches(5.0))
    s.shapes.add_picture(str(ASSETS / "phone_dashboard.png"), Inches(9.4), Inches(1.8), height=Inches(5.0))

    captions = [
        (1.0, "Feed", "Sub-second open, premium overlays, native tipping, watch-party."),
        (5.2, "Editor", "Multi-track, AI captions in 12s, BG removal, scene-cuts."),
        (9.4, "Dashboard", "Real-time earnings, audience demos, best-post-time regression."),
    ]
    for x, label, body in captions:
        add_text(s, x - 0.2, 6.85, 3.4, 0.3, label, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 5, total)

    # ============ SLIDE 6 — MARKET SIZE ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "05 · Market", "$300B TAM. $60B SAM. We need 0.5% to be a unicorn.")

    s.shapes.add_picture(str(ASSETS / "tam_sam_som.png"), Inches(0.5), Inches(1.9), width=Inches(8.0))

    # Right side narrative
    add_filled(s, 8.8, 1.9, 4.1, 4.7, WHITE)
    add_text(s, 9.0, 2.0, 3.8, 0.5, "How we scale into it", size=14, color=NAVY, bold=True)
    bullets = [
        ("Year 1", "iOS launch, EN markets, niche-down to one creator vertical (lifestyle / cooking)."),
        ("Year 2", "Android + Web, Spanish + Portuguese localisation. Brand marketplace GA."),
        ("Year 3", "International — Europe + LATAM. Creator Fund 2.0 at scale, ~$100M paid out."),
    ]
    for i, (y_label, body) in enumerate(bullets):
        y = 2.6 + i * 1.25
        add_filled(s, 9.0, y, 0.8, 0.4, PINK)
        add_text(s, 9.05, y + 0.04, 0.7, 0.4, y_label, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 9.0, y + 0.45, 3.8, 0.7, body, size=11, color=GRAY)

    add_footer(s, 6, total)

    # ============ SLIDE 7 — BUSINESS MODEL ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "06 · Business model", "Four revenue streams. Diversified by design.")

    s.shapes.add_picture(str(ASSETS / "unit_economics.png"), Inches(0.6), Inches(1.9), width=Inches(7.5))

    add_filled(s, 8.5, 1.9, 4.4, 4.7, WHITE)
    add_text(s, 8.7, 2.0, 4.0, 0.5, "Take rates", size=14, color=NAVY, bold=True)
    rows = [
        ("Tipping", "30%", "Stripe Connect, identifiable humans"),
        ("Subscriptions", "30% (Apple); 5% (web)", "Apple IAP for iOS, Stripe for web"),
        ("Brand deals", "10%", "Marketplace fee, escrow"),
        ("Premium content", "30% (Apple); 15% (web)", "Same routing as subs"),
    ]
    for i, (n, t, sub) in enumerate(rows):
        y = 2.6 + i * 0.95
        add_text(s, 8.7, y, 1.6, 0.3, n, size=11, color=NAVY, bold=True)
        add_text(s, 10.3, y, 2.6, 0.3, t, size=11, color=PINK, bold=True)
        add_text(s, 8.7, y + 0.32, 4.0, 0.5, sub, size=9, color=GRAY)

    add_footer(s, 7, total)

    # ============ SLIDE 8 — TRACTION + ROADMAP ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "07 · Traction & roadmap", "Code-complete. 6 architectural phases shipped. CI green. Month 0 = today.")

    s.shapes.add_picture(str(ASSETS / "traction_roadmap.png"), Inches(0.5), Inches(1.9), width=Inches(12.3))

    add_filled(s, 0.5, 6.0, 12.3, 0.85, WHITE)
    chips = [
        ("Phases 0-6 ✓", PINK),
        ("498 unit tests ✓", ACCENT),
        ("import-linter ✓", GOLD),
        ("FFmpeg pipeline ✓", NAVY),
        ("Stripe + IAP ✓", PINK),
        ("Live + duets ✓", ACCENT),
    ]
    cw = 12.0 / len(chips)
    for i, (label, c) in enumerate(chips):
        add_pill(s, 0.65 + i * cw, 6.18, cw - 0.15, 0.45, label, c, WHITE, size=11)

    add_footer(s, 8, total)

    # ============ SLIDE 9 — COMPETITION ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, LIGHT)
    title_block(s, "08 · Competition", "Pro tools live in editor apps. Creator economics live in walled gardens. We bring them together.")

    s.shapes.add_picture(str(ASSETS / "competitive_matrix.png"), Inches(0.5), Inches(1.6), height=Inches(5.4))

    add_filled(s, 7.5, 1.9, 5.4, 4.7, WHITE)
    add_text(s, 7.7, 2.0, 5.0, 0.5, "What's defensible", size=14, color=NAVY, bold=True)
    moats = [
        ("Architecture lead", "Clean DDD + 5 layer contracts. Adding creator features faster than incumbent monoliths."),
        ("Editor + feed coupling", "Pro tools pre-installed. CapCut + TikTok with one auth, one wallet, one analytic surface."),
        ("Transparent algo", "First major platform to expose Discovery Score features. Wins creator trust = wins creators."),
        ("Solo execution proof", "9–10mo to feature-complete vs. typical 18–24mo. Capital efficient."),
    ]
    for i, (h, b) in enumerate(moats):
        y = 2.6 + i * 1.05
        add_filled(s, 7.7, y, 0.18, 0.95, PINK)
        add_text(s, 8.0, y, 4.7, 0.4, h, size=12, color=NAVY, bold=True)
        add_text(s, 8.0, y + 0.4, 4.7, 0.6, b, size=10, color=GRAY)

    add_footer(s, 9, total)

    # ============ SLIDE 10 — THE ASK ============
    s = prs.slides.add_slide(blank)
    add_filled(s, 0, 0, 13.333, 7.5, DARK)
    add_filled(s, 0.5, 0.5, 0.15, 0.55, PINK)
    add_text(s, 0.8, 0.45, 8, 0.4, "09 · The ask", size=12, color=PINK, bold=True)
    add_text(s, 0.8, 0.75, 12, 0.7, "$3.5M seed for 18 months of runway and one verified revenue ramp.", size=24, color=WHITE, bold=True)

    s.shapes.add_picture(str(ASSETS / "use_of_funds.png"), Inches(0.8), Inches(1.9), width=Inches(7.0))

    # Right side terms
    add_filled(s, 8.0, 1.9, 4.9, 4.7, RGBColor(0x18, 0x18, 0x22))
    add_text(s, 8.2, 2.0, 4.6, 0.45, "Round terms", size=14, color=PINK, bold=True)

    terms = [
        ("Raise", "$3.5M"),
        ("Instrument", "SAFE, post-money cap $20M"),
        ("Lead check", "$1.0–1.5M"),
        ("Use period", "18 months"),
        ("Next milestone", "1M MAU @ month 12"),
        ("Series A trigger", "$3M ARR"),
        ("Investor seats", "1 board observer"),
        ("Founder commits", "100% time, 4-yr vest, 1-yr cliff"),
    ]
    for i, (k, v) in enumerate(terms):
        y = 2.55 + i * 0.42
        add_text(s, 8.2, y, 2.0, 0.3, k, size=11, color=RGBColor(0x9F, 0xAB, 0xBE))
        add_text(s, 10.2, y, 2.6, 0.3, v, size=11, color=WHITE, bold=True)

    # CTA bar
    add_filled(s, 0.8, 6.85, 11.7, 0.5, PINK)
    add_text(s, 0.8, 6.9, 11.7, 0.4,
             "Founder · asmeyatsky@hotmail.com · clipsmith.app",
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    out = "/Users/allansmeyatsky/clipsmith/clipsmith_seed_deck.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    main()

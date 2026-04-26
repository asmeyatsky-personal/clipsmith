"""Generate the 10-slide Clipsmith assessment deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x29, 0xB6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x5F, 0x6D)


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    s.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    s.shapes.title.text_frame.paragraphs[0].font.bold = True
    return s


def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]  # Title + content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title
    s.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    s.shapes.title.text_frame.paragraphs[0].font.bold = True

    body = s.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        if isinstance(b, tuple):
            label, detail = b
            p.text = ""
            r = p.add_run()
            r.text = label
            r.font.bold = True
            r.font.size = Pt(18)
            r.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = "  " + detail
            r2.font.size = Pt(16)
            r2.font.color.rgb = GRAY
        else:
            p.text = b
            for r in p.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = NAVY
        p.level = 0
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    add_title_slide(
        prs,
        "Clipsmith",
        "Solo-bootstrapped TikTok-class social video platform — engineering assessment",
    )

    # 2. The brief
    add_content_slide(prs, "The brief", [
        ("Product:", "Full-feature social video — feed, record, edit, monetize, live, AR."),
        ("Constraints:", "Solo engineer. Bootstrapped. iOS App Store target. Flexible timeline."),
        ("Method:", "9–10 month plan, 6 architectural phases, written into refactor250426.md."),
        ("Outcome:", "Phases 0–6 shipped, all CI green, ready for Phase 1.5 (TestFlight + submission)."),
    ])

    # 3. Architecture
    add_content_slide(prs, "Architecture — clean DDD with port/adapter discipline", [
        ("Layers:", "domain ← application ← infrastructure ← presentation."),
        ("Enforced:", "5 import-linter contracts in CI block layer crossings."),
        ("Stack:", "FastAPI · SQLModel · Alembic · Postgres (Neon) · Redis · RQ · Capacitor + Next.js."),
        ("Composition root:", "presentation/dependencies.py — single DI surface, exempt by design."),
        ("Code:", "~30k LOC backend Python, ~14k LOC frontend TypeScript, 90 SQL tables."),
    ])

    # 4. Phase 0–1: foundation + core
    add_content_slide(prs, "Phase 0–1 — foundation + social MVP", [
        ("Phase 0:", "Layer contracts, gitleaks, pip-audit, bandit, OpenTelemetry, Sentry, JSON logs."),
        ("Auth:", "Argon2id + JWT, 2FA via TOTP, email verify, COPPA-compliant age gate."),
        ("Social:", "Feed (foryou/following/trending), follows, blocks, likes, comments, hashtags."),
        ("Compliance:", "GDPR data export + delete, audit log, content moderation port + OpenAI adapter."),
        ("Frontend:", "Next.js static export, Capacitor iOS shell, mobile-first feed."),
    ])

    # 5. Phase 2–3: AI + monetization
    add_content_slide(prs, "Phase 2–3 — AI suite + monetization", [
        ("AI suite:", "Captions (AssemblyAI port), scene detect (PySceneDetect), voice enhance (RNNoise)."),
        ("Editor:", "Multi-track timeline, BG removal (MediaPipe), 10 starter templates, AI-driven."),
        ("Monetization:", "Stripe Connect tipping, premium-gated content, subscription tiers."),
        ("Creator Fund:", "70/30 split, monthly payout batch task, ledger + wallet entities."),
        ("Brand marketplace:", "Profile, campaign, response endpoints scaffolded."),
    ])

    # 6. Phase 4–5: live, community, AR
    add_content_slide(prs, "Phase 4–5 — live, community, AR", [
        ("Duets:", "FFmpeg side-by-side / vstack / PiP composition, queued automatically."),
        ("Watch parties:", "WebSocket room broadcasts play/pause/seek with JWT auth."),
        ("Live streaming:", "LiveStreamPort + placeholder adapter (Mediasoup deferred to VPS day-1)."),
        ("Community:", "Groups, posts, events, RSVPs, circles, polls, chapters, badges, challenges."),
        ("AR + grading:", "Effects marketplace (8 starters), 6 LUT presets, server-side chroma-key."),
    ])

    # 7. Phase 6 + non-functional
    add_content_slide(prs, "Phase 6 — i18n, scale, polish", [
        ("i18n:", "EN, ES, PT, FR, DE, JA, KO bundles + locale-aware <LanguageSwitcher />."),
        ("Performance:", "Feed weak-ETag → 304 short-circuit + Cache-Control: private, max-age=10."),
        ("Analytics export:", "AnalyticsExportPort + JSONL adapter; daily snapshot task."),
        ("Discovery:", "Real per-video score (virality/freshness/community/interest) — transparent."),
        ("Posting recs:", "Regression on creator's own video history, weekday × hour buckets."),
    ])

    # 8. Quality
    add_content_slide(prs, "Quality posture", [
        ("Tests:", "498 unit + 10 e2e smoke + 5 layer contracts — all green on every push."),
        ("CI:", "GitHub Actions: secrets-scan, lint-imports, pytest+cov≥60, smoke, bandit, pip-audit, npm-audit≥high, docker."),
        ("Security:", "All known high/critical CVEs resolved (next, dotenv, jose, cryptography, …)."),
        ("Observability:", "Correlation IDs, structlog JSON, OTel hooks, Sentry, Prometheus."),
        ("Migrations:", "Alembic single-head; tested against Postgres on Neon."),
    ])

    # 9. What's left
    add_content_slide(prs, "What is left after Phase 6", [
        ("Phase 1.5 — TestFlight + App Store:", "Apple Developer setup, IAP wiring, signing, screenshots, App Review (~3 weeks)."),
        ("Deferred adapters:", "Real Mediasoup SFU when VPS is provisioned. Real BigQuery exporter."),
        ("Conditional:", "Rust hot-path migration only if measured p99 > 50ms."),
        ("Non-engineering risks:", "Music licensing strategy, mod contractor at DAU > 500, GDPR DPAs, distribution wedge."),
        ("Out of scope:", "External pen test (do at DAU > 1k or after first incident)."),
    ])

    # 10. Risks + ask
    add_content_slide(prs, "Top risks and decisions you own", [
        ("App Review cycles:", "Budget 2–3 rejection rounds for a UGC + payments app."),
        ("Music rights:", "PRD allocated $5M/yr; we have $0. Royalty-free libraries only at launch."),
        ("Cold-start network effects:", "Engineering doesn't solve this — niche down to one creator vertical."),
        ("Solo burnout:", "Highest-probability failure mode. Take time off between phases."),
        ("Founder decisions due:", "App name, IAP pricing, primary domain, marketing-day plan."),
    ])

    out = "/Users/allansmeyatsky/clipsmith/clipsmith_assessment.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    main()
